from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from src.config import ConfigOCR, Settings, get_config_technique, get_settings

logger = logging.getLogger(__name__)


# ===========================================================================
# Structures de retour
# ===========================================================================

@dataclass
class Page:
    """Une page (ou une feuille de tableur, ou un bloc pour les formats sans pagination)."""
    numero: int
    texte: str


@dataclass
class DocumentCharge:
    """Résultat de l'extraction d'un fichier."""
    chemin: Path
    pages: list[Page]
    ocr_utilise: bool = False
    avertissements: list[str] = field(default_factory=list)

    @property
    def texte_complet(self) -> str:
        return "\n\n".join(p.texte for p in self.pages)

    @property
    def nb_caracteres(self) -> int:
        return sum(len(p.texte) for p in self.pages)

    @property
    def est_vide(self) -> bool:
        return self.nb_caracteres == 0


class ErreurChargement(Exception):
    """Levée quand un fichier ne peut être lu par aucun loader disponible."""


# ===========================================================================
# Loaders par format
# ===========================================================================

def _charger_texte(chemin: Path) -> list[Page]:
    """.txt / .md : lecture directe, un seul bloc."""
    contenu = chemin.read_text(encoding="utf-8", errors="replace")
    return [Page(numero=1, texte=contenu.strip())]


def _table_en_lignes(table: list[list[str | None]]) -> list[str]:
    """
    Convertit une table pdfplumber en lignes « a | b | c ».

    Le format est identique à celui produit par les loaders docx, xlsx et
    csv : le découpage structure-aware reconnaît ainsi les tableaux de tous
    les formats avec la même règle.
    """
    lignes: list[str] = []
    for ligne in table or []:
        cellules = [
            " ".join(str(cellule).split()) if cellule is not None else ""
            for cellule in ligne
        ]
        if any(cellules):
            lignes.append(" | ".join(cellules))
    return lignes


def _tables_fiables(page: Any) -> list[list[str]]:
    """
    Extrait les tableaux d'une page, en écartant ceux qui n'en sont pas.

    pdfplumber détecte parfois des « tableaux » d'une seule colonne sur du
    texte mis en page. Un tableau est retenu s'il a au moins deux lignes et
    deux colonnes exploitables.
    """
    try:
        tables = page.extract_tables() or []
    except Exception as exc:  # noqa: BLE001 — extraction best effort
        logger.debug("Extraction de tableaux impossible : %s", exc)
        return []

    retenues: list[list[str]] = []
    for table in tables:
        lignes = _table_en_lignes(table)
        if len(lignes) >= 2 and all(ligne.count("|") >= 1 for ligne in lignes[:2]):
            retenues.append(lignes)
    return retenues


def _sans_doublons_de_table(texte: str, tables: list[list[str]]) -> str:
    """
    Retire du texte brut les lignes déjà couvertes par un tableau extrait.

    Sans cette étape, chaque valeur d'un tableau apparaîtrait deux fois dans
    la page : une fois en texte désordonné, une fois en ligne structurée.
    La comparaison se fait sur les cellules, car pdfplumber restitue les
    mêmes valeurs avec des espacements différents selon la méthode.
    """
    if not tables:
        return texte

    empreintes = {
        _empreinte_ligne(ligne) for lignes in tables for ligne in lignes
    }
    conservees = [
        ligne
        for ligne in texte.splitlines()
        if not ligne.strip() or _empreinte_ligne(ligne) not in empreintes
    ]
    return "\n".join(conservees)


def _empreinte_ligne(ligne: str) -> str:
    """Signature d'une ligne, insensible aux espaces et aux séparateurs."""
    return " ".join(ligne.replace("|", " ").split()).casefold()


def _charger_pdf(chemin: Path) -> list[Page]:
    """
    PDF : extraction de la couche texte via pdfplumber.

    pdfplumber gère mieux la mise en page multi-colonne que pypdf.
    Les tableaux détectés sont extraits séparément et rendus au format
    « a | b | c », puis retirés du texte brut pour éviter de dupliquer les
    mêmes valeurs. Si aucun tableau fiable n'est détecté, le comportement
    se réduit à l'ancien `extract_text()`.
    L'OCR éventuel est décidé plus haut, pas ici.
    """
    import pdfplumber

    pages: list[Page] = []
    with pdfplumber.open(chemin) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            texte = page.extract_text() or ""
            tables = _tables_fiables(page)

            if tables:
                texte = _sans_doublons_de_table(texte, tables)
                blocs_tables = ["\n".join(lignes) for lignes in tables]
                texte = "\n\n".join([texte.strip(), *blocs_tables]).strip()

            pages.append(Page(numero=i, texte=texte.strip()))
    return pages


def _charger_docx(chemin: Path) -> list[Page]:
    """
    DOCX : paragraphes + tableaux. Word n'a pas de pagination fiable
    hors rendu, donc tout est regroupé en un bloc.
    """
    from docx import Document

    doc = Document(str(chemin))
    morceaux: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]

    for table in doc.tables:
        for ligne in table.rows:
            cellules = [c.text.strip() for c in ligne.cells]
            if any(cellules):
                morceaux.append(" | ".join(cellules))

    return [Page(numero=1, texte="\n".join(morceaux).strip())]


def _charger_xlsx(chemin: Path) -> list[Page]:
    """
    XLSX : une Page par feuille, chaque ligne en Markdown ' | '.
    Le format tabulaire est préservé pour rester lisible par le LLM.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(chemin), read_only=True, data_only=True)
    pages: list[Page] = []

    for idx, feuille in enumerate(wb.worksheets, start=1):
        lignes: list[str] = []
        for ligne in feuille.iter_rows(values_only=True):
            cellules = [str(c) if c is not None else "" for c in ligne]
            if any(cellules):
                lignes.append(" | ".join(cellules))
        if lignes:
            entete = f"# Feuille : {feuille.title}"
            pages.append(Page(numero=idx, texte=entete + "\n" + "\n".join(lignes)))

    wb.close()
    return pages or [Page(numero=1, texte="")]


def _charger_csv(chemin: Path) -> list[Page]:
    """CSV : lecture via pandas, rendu Markdown."""
    import pandas as pd

    try:
        df = pd.read_csv(chemin, dtype=str, keep_default_na=False)
    except UnicodeDecodeError:
        df = pd.read_csv(chemin, dtype=str, keep_default_na=False, encoding="latin-1")

    lignes = [" | ".join(df.columns)]
    for _, row in df.iterrows():
        lignes.append(" | ".join(str(v) for v in row.values))
    return [Page(numero=1, texte="\n".join(lignes))]


def _charger_pptx(chemin: Path) -> list[Page]:
    """PPTX : une Page par diapositive."""
    from pptx import Presentation

    prs = Presentation(str(chemin))
    pages: list[Page] = []
    for i, slide in enumerate(prs.slides, start=1):
        morceaux = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        pages.append(Page(numero=i, texte="\n".join(morceaux)))
    return pages


def _charger_html(chemin: Path) -> list[Page]:
    """HTML : texte visible seulement, scripts et styles retirés."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(chemin.read_text(encoding="utf-8", errors="replace"), "html.parser")
    for balise in soup(["script", "style", "noscript"]):
        balise.decompose()
    texte = soup.get_text(separator="\n")
    lignes = [l.strip() for l in texte.splitlines() if l.strip()]
    return [Page(numero=1, texte="\n".join(lignes))]


# Table extension -> loader. Ajouter un format = ajouter une ligne.
_LOADERS = {
    ".txt": _charger_texte,
    ".md": _charger_texte,
    ".pdf": _charger_pdf,
    ".docx": _charger_docx,
    ".xlsx": _charger_xlsx,
    ".csv": _charger_csv,
    ".pptx": _charger_pptx,
    ".html": _charger_html,
    ".htm": _charger_html,
}


# ===========================================================================
# OCR (PDF scannés)
# ===========================================================================

def _configurer_tesseract(settings: Settings) -> None:
    """Renseigne le chemin du binaire si fourni (nécessaire sous Windows)."""
    if settings.tesseract_cmd:
        import pytesseract
        pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd


def _ocr_pdf(chemin: Path, cfg_ocr: ConfigOCR, settings: Settings) -> list[Page]:
    """
    Rendu de chaque page en image puis OCR Tesseract, PAR LOTS de
    `cfg_ocr.pages_max` pages (ex. 300 pages, pages_max=20 -> 15 lots
    successifs 1-20, 21-40, ..., 281-300, fusionnés dans l'ordre).

    `pages_max` borne désormais le coût (mémoire, durée) d'UN lot de rendu,
    plus jamais le document entier : l'ancien comportement rendait TOUTES
    les pages puis ne conservait que les `pages_max` premières, abandonnant
    silencieusement le reste (voir l'audit long-documents — un PDF scanné
    de 300 pages n'était en réalité traité que sur ses 20 premières, sans
    qu'aucun appelant ne puisse le détecter). Un document scanné, quelle
    que soit sa longueur, est désormais intégralement couvert.
    """
    import pytesseract
    from pdf2image import convert_from_path, pdfinfo_from_path

    _configurer_tesseract(settings)

    try:
        nb_pages = int(pdfinfo_from_path(str(chemin)).get("Pages") or 0)
    except Exception as exc:  # noqa: BLE001 — repli : rendu en un seul lot, sans troncature
        logger.warning(
            "%s : nombre de pages indisponible (%s), OCR en un seul lot.",
            chemin.name, exc,
        )
        nb_pages = 0

    pas = max(cfg_ocr.pages_max, 1)
    bornes = (
        [(debut, min(debut + pas - 1, nb_pages)) for debut in range(1, nb_pages + 1, pas)]
        if nb_pages > 0
        else [(None, None)]  # nombre de pages inconnu : rendu complet, un seul lot
    )

    pages: list[Page] = []
    for debut, fin in bornes:
        if debut is not None:
            logger.info("%s : OCR pages %d-%d/%d…", chemin.name, debut, fin, nb_pages)
        images = convert_from_path(
            str(chemin), dpi=cfg_ocr.dpi, first_page=debut, last_page=fin,
        )
        decalage = (debut - 1) if debut is not None else 0
        for i, image in enumerate(images, start=1):
            texte = pytesseract.image_to_string(image, lang=cfg_ocr.langues)
            pages.append(Page(numero=decalage + i, texte=texte.strip()))

    return pages


# ===========================================================================
# Point d'entrée
# ===========================================================================

def charger_document(chemin: Path) -> DocumentCharge:
    """
    Charge un fichier et renvoie son texte, page par page.

    Bascule sur l'OCR si un PDF rend un texte anormalement court
    (couche texte absente = scan). Lève ErreurChargement si le format
    n'est pas géré ou si l'extraction échoue totalement.
    """
    settings = get_settings()
    tech = get_config_technique()
    ext = chemin.suffix.lower()

    loader = _LOADERS.get(ext)
    if loader is None:
        raise ErreurChargement(f"Format non géré : {ext} ({chemin.name})")

    try:
        pages = loader(chemin)
    except Exception as exc:  # noqa: BLE001 — on requalifie en erreur métier
        raise ErreurChargement(f"Échec d'extraction sur {chemin.name} : {exc}") from exc

    doc = DocumentCharge(chemin=chemin, pages=pages)

    # Détection de scan : uniquement pour les PDF, si OCR activé.
    scan_probable = (
        ext == ".pdf"
        and tech.ocr.active
        and settings.ocr_enabled
        and doc.nb_caracteres < tech.ingestion.seuil_texte_vide
    )

    if scan_probable:
        logger.info("%s : texte court (%d car.), OCR…", chemin.name, doc.nb_caracteres)
        try:
            pages_ocr = _ocr_pdf(chemin, tech.ocr, settings)
            doc = DocumentCharge(chemin=chemin, pages=pages_ocr, ocr_utilise=True)
            if doc.est_vide:
                doc.avertissements.append("OCR effectué mais aucun texte récupéré.")
        except Exception as exc:  # noqa: BLE001
            logger.error("OCR échoué sur %s : %s", chemin.name, exc)
            doc.avertissements.append(f"OCR indisponible ou échoué : {exc}")

    if doc.est_vide:
        doc.avertissements.append("Aucun texte extrait — document ignoré à l'indexation.")

    return doc


def formats_supportes() -> list[str]:
    return sorted(_LOADERS.keys())